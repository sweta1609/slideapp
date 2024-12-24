
from flask import Flask, jsonify
from pptx import Presentation

app = Flask(__name__)

@app.route('/generate_presentation', methods=['GET'])
def generate_presentation():
    try:
        # Create a presentation and add a simple slide
        presentation = Presentation()
        slide_layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Generated Presentation Slide"

        # Save the presentation
        presentation.save('generated_presentation.pptx')

        return jsonify({'status': 'Presentation generated successfully'})
    except Exception as e:
        return jsonify({'status': f'Error: {str(e)}'}), 500

if __name__ == '__main__':
    app.run(debug=True)
