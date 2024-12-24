from flask import Flask, request, jsonify, send_file
from pptx import Presentation
from io import BytesIO

app = Flask(__name__)

@app.route('/edit_presentation', methods=['POST'])
def edit_presentation():
    try:
        # Read presentation from the request
        presentation_file = request.files['presentation']
        presentation = Presentation(BytesIO(presentation_file.read()))

        # Example: Add a new slide
        slide_layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Edited Slide"

        # Save the edited presentation to a BytesIO buffer
        edited_presentation_buffer = BytesIO()
        presentation.save(edited_presentation_buffer)
        edited_presentation_buffer.seek(0)

        # Send the edited presentation as a response
        return send_file(
            edited_presentation_buffer,
            as_attachment=True,
            download_name='edited_presentation.pptx',
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
    except Exception as e:
        return jsonify({'status': f'Error: {str(e)}'}), 500

if __name__ == '__main__':
    app.run(debug=True)
